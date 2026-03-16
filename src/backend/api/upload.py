"""Upload API — Secure File Ingestion & Content Extraction Endpoints
====================================================================

Handles multi-format file uploads (CSV, JSON, PDF, DOCX, XLSX, PPTX, RTF,
Parquet, Feather, HDF5, NetCDF, MAT) with layered security validation:

1. **Filename sanitization** — path-traversal prevention via werkzeug.
2. **Extension allow-list** — rejects unknown file types.
3. **Size enforcement** — streaming chunk check against ``MAX_FILE_SIZE``.
4. **MIME/content validation** — ``python-magic`` or signature-based fallback.
5. **Text sanitization** — bleach / HTML-escape before storage.
6. **Secure permissions** — ``0600`` on saved files.

After upload, structured files are column-scanned, text documents are
chunked and indexed into ChromaDB for RAG retrieval, and the analysis
cache is tag-invalidated so stale results are never served.

Endpoints
---------
``POST /``
    Upload a single file with full security pipeline.
``POST /raw-text``
    Accept raw text input without file upload.
``GET  /files``
    List all uploaded files with metadata.
``GET  /files/{filename}/summary``
    Return column & statistical summary for a structured file.
``DELETE /files/{filename}``
    Remove an uploaded file and its extracted text.
``GET  /templates``
    Provide sample file download links for first-time users.
"""

from __future__ import annotations

from fastapi import APIRouter, UploadFile, File
import os
import tempfile
import shutil
import stat
import re
import datetime
import logging
from werkzeug.utils import secure_filename
from backend.utils.data_utils import read_dataframe, validate_dataframe, create_data_summary
try:
    import pypdf
except ImportError:
    pypdf = None

try:
    import openpyxl
    HAS_OPENPYXL = True
except ImportError:
    HAS_OPENPYXL = False
    openpyxl = None

try:
    from docx import Document
    HAS_PYTHON_DOCX = True
except ImportError:
    HAS_PYTHON_DOCX = False
    Document = None

try:
    from pptx import Presentation
    HAS_PYTHON_PPTX = True
except ImportError:
    HAS_PYTHON_PPTX = False
    Presentation = None

try:
    import chardet
    HAS_CHARDET = True
except ImportError:
    HAS_CHARDET = False
    chardet = None

try:
    import magic
    HAS_MAGIC = True
except ImportError:
    HAS_MAGIC = False
    magic = None

try:
    import filetype
    HAS_FILETYPE = True
except ImportError:
    HAS_FILETYPE = False
    filetype = None

try:
    import pandas as pd
except ImportError:
    pd = None

try:
    import bleach
    HAS_BLEACH = True
except ImportError:
    HAS_BLEACH = False
    bleach = None

from pydantic import BaseModel

logger = logging.getLogger(__name__)
router = APIRouter()


class RawTextInput(BaseModel):
    """Request body for direct text ingestion without file upload.

    Attributes:
        text:        The raw text content to analyze.
        title:       Human-readable title (used to generate the saved filename).
        description: Optional context about the text content.
    """

    text: str
    title: str = "Raw Text Input"
    description: str = ""

# Security Configuration
MAX_FILE_SIZE: int = 100 * 1024 * 1024  # 100 MB upload ceiling
# Phase 3.5: Added scientific file formats (parquet, feather, hdf5, h5, nc, mat)
ALLOWED_EXTENSIONS = {
    # Standard formats
    '.csv', '.json', '.pdf', '.txt', '.xlsx', '.xls', '.docx', '.pptx', '.rtf',
    # Scientific formats (Phase 3.5)
    '.parquet', '.feather', '.hdf5', '.h5', '.nc', '.mat'
}
ALLOWED_MIME_TYPES = {
    '.csv': ['text/csv', 'text/plain', 'application/csv'],
    '.json': ['application/json', 'text/plain'],
    '.pdf': ['application/pdf'],
    '.txt': ['text/plain'],
    '.xlsx': ['application/vnd.openxmlformats-officedocument.spreadsheetml.sheet'],
    '.xls': ['application/vnd.ms-excel'],
    '.docx': ['application/vnd.openxmlformats-officedocument.wordprocessingml.document'],
    '.pptx': ['application/vnd.openxmlformats-officedocument.presentationml.presentation'],
    '.rtf': ['application/rtf', 'text/rtf'],
    # Phase 3.5: Scientific formats MIME types
    '.parquet': ['application/octet-stream', 'application/x-parquet'],
    '.feather': ['application/octet-stream', 'application/x-feather'],
    '.hdf5': ['application/x-hdf5', 'application/octet-stream'],
    '.h5': ['application/x-hdf5', 'application/octet-stream'],
    '.nc': ['application/x-netcdf', 'application/octet-stream'],
    '.mat': ['application/x-matlab-data', 'application/octet-stream']
}

# Use centralized path resolver
from backend.utils.data_utils import DataPathResolver
DataPathResolver.ensure_directories_exist()

def validate_filename(filename: str) -> str:
    """Sanitize and validate a user-supplied filename against path-traversal attacks.

    Args:
        filename: Raw filename from the upload request.

    Returns:
        A sanitized filename safe for use on the local filesystem.

    Raises:
        ValueError: If the filename is missing, contains traversal sequences,
            null bytes, or becomes empty after sanitization.
    """
    if not filename:
        raise ValueError("Filename is required")
    
    if filename is None:
        raise ValueError("Filename cannot be None")
    
    # Check for obvious path traversal attempts before secure_filename
    if '..' in filename or '/' in filename or '\\' in filename:
        raise ValueError("Path traversal attempt detected in filename")
    
    # Check for null bytes and other dangerous characters
    if '\x00' in filename or any(ord(c) < 32 for c in filename):
        raise ValueError("Invalid characters detected in filename")
    
    # Use secure_filename for additional sanitization
    filename = secure_filename(filename)
    if not filename:
        raise ValueError("Filename becomes empty after sanitization")
    
    # Final validation - must have valid extension and reasonable length
    if len(filename) > 255:
        raise ValueError("Filename too long")
    
    if '.' not in filename:
        raise ValueError("Filename must have an extension")
    
    return filename

def validate_file_size(size: int) -> None:
    """Raise ``ValueError`` if *size* exceeds ``MAX_FILE_SIZE``."""
    if size > MAX_FILE_SIZE:
        raise ValueError(f"File size {size} exceeds maximum allowed size of {MAX_FILE_SIZE} bytes ({MAX_FILE_SIZE // (1024*1024)}MB)")

def validate_file_extension(filename: str) -> str:
    """Return the lowercase extension if it is in ``ALLOWED_EXTENSIONS``; raise otherwise."""
    ext = os.path.splitext(filename)[1].lower()
    if ext not in ALLOWED_EXTENSIONS:
        raise ValueError(f"File extension '{ext}' not allowed. Supported extensions: {', '.join(ALLOWED_EXTENSIONS)}")
    return ext

def validate_file_content(content: bytes, extension: str) -> bool:
    """Verify file content matches the claimed extension via MIME detection.

    Uses ``python-magic`` (libmagic) when available, falling back to the
    ``filetype`` package and then lightweight signature checks.

    Args:
        content:   First 1 KB of the uploaded file.
        extension: Lowercase extension including the leading dot.

    Returns:
        ``True`` if the content is consistent with the extension.
    """
    # Prefer python-magic (libmagic) if available for robust detection
    if HAS_MAGIC:
        try:
            mime_type = magic.from_buffer(content[:1024], mime=True)  # Check first 1KB
            allowed_mimes = ALLOWED_MIME_TYPES.get(extension, [])

            is_valid = mime_type in allowed_mimes
            if not is_valid:
                logger.warning("MIME type mismatch: got %s, expected one of %s for extension %s", mime_type, allowed_mimes, extension)

            return is_valid
        except Exception as e:
            logger.error("Failed to validate file content with python-magic: %s", e, exc_info=True)
            return False

    # Fallback: try the 'filetype' package if available
    logger.info("python-magic not available, using fallback content checks. For best results on Windows, consider installing python-magic-bin or libmagic.")
    try:
        if HAS_FILETYPE:
            try:
                kind = filetype.guess(content)
                mime_type = getattr(kind, 'mime', None)
                if mime_type:
                    allowed_mimes = ALLOWED_MIME_TYPES.get(extension, [])
                    if mime_type in allowed_mimes:
                        return True
                    else:
                        logger.warning("Fallback filetype mismatch: got %s, expected one of %s for extension %s", mime_type, allowed_mimes, extension)
                        return False
            except Exception as e:
                logger.debug("filetype.guess failed: %s", e)

        # Lightweight signature-based checks
        # PDF files start with %PDF
        if extension == '.pdf':
            return content.startswith(b'%PDF')

        # OOXML and similar formats (xlsx, docx, pptx) are ZIP packages starting with PK
        if extension in ('.xlsx', '.xls', '.docx', '.pptx'):
            return content.startswith(b'PK')

        # Text-like formats: attempt to decode as utf-8 or latin-1
        if extension in ('.txt', '.csv', '.json', '.rtf'):
            try:
                content.decode('utf-8')
                return True
            except Exception:
                try:
                    content.decode('latin-1')
                    return True
                except Exception:
                    logger.warning("Unable to decode text-like file for extension %s", extension)
                    return False

        # Unknown extension: be permissive but log the uncertainty
        logger.info("Fallback content validation could not determine MIME reliably; allowing upload as a permissive fallback")
        return True
    except Exception as e:
        logger.error("Unexpected error during fallback content validation: %s", e, exc_info=True)
        return False

def sanitize_extracted_text(text: str) -> str:
    """Strip HTML tags and truncate extracted text to prevent XSS and memory issues."""
    if not text:
        return ""
    
    # Limit text length to prevent memory issues
    max_text_length = 1024 * 1024  # 1MB of text
    if len(text) > max_text_length:
        text = text[:max_text_length]
        logger.warning("Truncated extracted text to %s characters", max_text_length)
    
    # Use bleach for HTML sanitization if available
    if HAS_BLEACH:
        # Remove all HTML tags and dangerous content
        text = bleach.clean(text, tags=[], attributes={}, strip=True)
    else:
        # Basic HTML escaping if bleach not available
        import html
        text = html.escape(text)
    
    return text

def secure_file_path(filename: str) -> str:
    """Build an absolute path within the uploads directory and verify no traversal.

    Raises:
        ValueError: If the resolved path escapes the uploads directory.
    """
    data_dir = str(DataPathResolver.get_uploads_dir())
    file_path = os.path.join(data_dir, filename)
    
    # Resolve path to prevent directory traversal
    resolved_path = os.path.realpath(file_path)
    resolved_data_dir = os.path.realpath(data_dir)
    
    # Ensure the file path is within DATA_DIR
    if not resolved_path.startswith(resolved_data_dir):
        raise ValueError(f"Path traversal attempt detected: {resolved_path}")
    
    return resolved_path


# Secure file upload endpoint with comprehensive security validation
@router.post("/")
async def upload_document(file: UploadFile = File(...)):
    """Secure file upload with multi-format support and layered validation.

    Accepts CSV, JSON, PDF, TXT, XLSX/XLS, DOCX, PPTX, RTF, Parquet,
    Feather, HDF5, NetCDF, and MAT files.  Each upload passes through
    filename sanitisation, extension allow-list, streaming size check,
    MIME content validation, text extraction, ChromaDB indexing, and
    cache invalidation.

    Args:
        file: The uploaded file from the multipart form.

    Returns:
        Dict with ``filename``, ``message``, ``file_size``, optional
        ``columns``, and ``extracted_text_path`` on success; ``error``
        on failure.
    """
    temp_file_path = None
    final_file_path = None
    
    try:
        logger.debug("[UPLOAD] Received upload request for file: %s", file.filename)
        
        # Step 1: Validate filename
        try:
            filename = validate_filename(file.filename)
        except ValueError as e:
            logger.warning("[UPLOAD] Filename validation failed: %s", e)
            return {"error": f"Invalid filename: {str(e)}"}
        
        # Step 2: Validate file extension  
        try:
            extension = validate_file_extension(filename)
        except ValueError as e:
            logger.warning("[UPLOAD] File extension validation failed: %s", e)
            return {"error": str(e)}
        
        # Step 3: Check file size if available
        if file.size:
            try:
                validate_file_size(file.size)
            except ValueError as e:
                logger.warning("[UPLOAD] File size validation failed: %s", e)
                return {"error": str(e)}
        
        # Step 4: Create secure file path
        try:
            final_file_path = secure_file_path(filename)
        except ValueError as e:
            logger.error("[UPLOAD] Path validation failed: %s", e, exc_info=True)
            return {"error": "Invalid file path"}
        
        # Step 5: Process file upload with temporary file for security  
        temp_file_path = None
        try:
            with tempfile.NamedTemporaryFile(delete=False) as temp_file:
                temp_file_path = temp_file.name
                
                # Stream file in chunks to prevent memory exhaustion
                total_size = 0
                chunk_size = 8192  # 8KB chunks
                
                while chunk := await file.read(chunk_size):
                    chunk_size_bytes = len(chunk)
                    total_size += chunk_size_bytes
                    
                    # Enforce size limit during streaming
                    if total_size > MAX_FILE_SIZE:
                        raise ValueError(f"File size exceeds maximum limit of {MAX_FILE_SIZE // (1024*1024)}MB during upload")
                    
                    temp_file.write(chunk)
                
                temp_file.flush()
                logger.debug("[UPLOAD] Successfully streamed %s bytes to temporary file", total_size)
            
            # File is now closed and can be moved safely
            
            # Step 6: Validate file content matches extension
            with open(temp_file_path, 'rb') as f:
                content_sample = f.read(1024)  # Read first 1KB for validation
                if not validate_file_content(content_sample, extension):
                    logger.warning("[UPLOAD] File content validation failed for %s", filename)
                    return {"error": "File content does not match the file extension"}
            
            # Step 7: Move file to final location with secure permissions
            shutil.move(temp_file_path, final_file_path)
            temp_file_path = None  # File has been moved, don't try to delete it
            
            # Set restrictive permissions (owner read/write only)
            os.chmod(final_file_path, stat.S_IRUSR | stat.S_IWUSR)
            logger.debug("[UPLOAD] File saved securely to: %s", final_file_path)
            
            # Step 8: Extract and sanitize text content if applicable
            extracted_text = None
            extracted_text_path = None
            
            if extension == '.pdf':
                extracted_text = await extract_pdf_text_secure(final_file_path, filename)
            elif extension == '.txt':
                extracted_text = await extract_txt_text_secure(final_file_path, filename)
            elif extension in ['.xlsx', '.xls']:
                extracted_text = await extract_excel_text_secure(final_file_path, filename)
            elif extension == '.docx':
                extracted_text = await extract_docx_text_secure(final_file_path, filename)
            elif extension == '.pptx':
                extracted_text = await extract_pptx_text_secure(final_file_path, filename)
            elif extension == '.rtf':
                extracted_text = await extract_rtf_text_secure(final_file_path, filename)
            
            # Step 9: Save sanitized extracted text if available
            if extracted_text:
                sanitized_text = sanitize_extracted_text(extracted_text)
                extracted_text_path = final_file_path + '.extracted.txt'
                
                with open(extracted_text_path, 'w', encoding='utf-8') as tf:
                    tf.write(sanitized_text)
                
                # Set secure permissions on extracted text file
                os.chmod(extracted_text_path, stat.S_IRUSR | stat.S_IWUSR)
                logger.debug("[UPLOAD] Sanitized extracted text saved to: %s", extracted_text_path)
                
                # Step 9.5: Index extracted text into ChromaDB for RAG retrieval
                try:
                    from backend.core.chromadb_client import ChromaDBClient, chunk_text
                    
                    chroma_client = ChromaDBClient()
                    chunks = chunk_text(sanitized_text, chunk_size=1000, overlap=200)
                    
                    for i, chunk in enumerate(chunks):
                        doc_id = f"{filename}_chunk_{i}"
                        chroma_client.add_document(
                            doc_id=doc_id,
                            text=chunk,
                            metadata={"filename": filename, "chunk_index": i, "total_chunks": len(chunks)}
                        )
                    
                    logger.debug("[UPLOAD] Indexed %s chunks from %s into ChromaDB", len(chunks), filename)
                except Exception as index_error:
                    logger.warning("[UPLOAD] Failed to index document into ChromaDB: %s", index_error)
                    # Don't fail upload if indexing fails
            
            # Step 10: Extract column information for structured files
            columns = None
            if extension == '.csv':
                columns = await extract_csv_columns(final_file_path, filename)
            elif extension in ['.xlsx', '.xls']:
                columns = await extract_excel_columns(final_file_path, filename)
            
            # Step 11a: Invalidate DataFrame store so stale in-memory copies are dropped
            try:
                from backend.core.dataframe_store import get_dataframe_store
                get_dataframe_store().invalidate(final_file_path)
                logger.debug("[UPLOAD] DataFrame store invalidated for: %s", filename)
            except Exception as store_error:
                logger.debug("[UPLOAD] DataFrame store invalidation skipped: %s", store_error)

            # Step 11b: Invalidate analysis cache for this filename to prevent stale data
            try:
                from backend.core.enhanced_cache_integration import get_enhanced_cache_manager
                cache_mgr = get_enhanced_cache_manager()
                
                # Tag-based invalidation across all tiers (L1/L2/L3)
                cache_mgr.l3_cache.invalidate_by_tags({filename, 'structured_data', 'file_analysis'})
                
                # Clear L1 (LRU) entries containing filename
                if hasattr(cache_mgr.l1_cache, 'cache'):
                    keys_to_remove = [k for k in list(cache_mgr.l1_cache.cache.keys()) if filename in str(k)]
                    for k in keys_to_remove:
                        cache_mgr.l1_cache.cache.pop(k, None)
                
                logger.debug("[UPLOAD] Cache invalidated for file: %s", filename)
            except Exception as cache_error:
                logger.debug("[UPLOAD] Cache invalidation skipped: %s", cache_error)
                # Don't fail the upload if cache invalidation fails
            
            logger.debug("[UPLOAD] Upload completed successfully for: %s", filename)
            return {
                "filename": filename,
                "message": "File uploaded successfully",
                "file_size": total_size,
                "extracted_text_path": extracted_text_path,
                "columns": columns
            }
                
        except ValueError as e:
            logger.error("[UPLOAD] Validation error during upload: %s", e, exc_info=True)
            return {"error": str(e)}
        
        except Exception as e:
            logger.error("[UPLOAD] Unexpected error during file processing: %s", e, exc_info=True)
            return {"error": "File processing failed due to an internal error"}
    
    except Exception as e:
        logger.error("[UPLOAD] Critical error in upload endpoint: %s", e, exc_info=True)
        return {"error": "Upload failed due to an internal error"}
    
    finally:
        # Cleanup: Remove temporary file if it still exists
        if temp_file_path and os.path.exists(temp_file_path):
            try:
                os.unlink(temp_file_path)
                logger.info("[UPLOAD] Cleaned up temporary file: %s", temp_file_path)
            except OSError as e:
                logger.warning("[UPLOAD] Failed to cleanup temporary file %s: %s", temp_file_path, e)
        
        # NOTE: Final file cleanup is intentionally skipped — the file is the user's upload
        # and should persist. Temporary files are cleaned above.

async def extract_pdf_text_secure(file_path: str, filename: str) -> str:
    """Securely extract text from PDF file with proper error handling"""
    if pypdf is None:
        logger.error("[UPLOAD] pypdf not available for PDF text extraction: %s", filename, exc_info=True)
        raise ValueError("PDF text extraction not available - pypdf not installed")
    
    try:
        with open(file_path, 'rb') as f:
            pdf_reader = pypdf.PdfReader(f)
            
            # Limit number of pages to prevent DoS
            max_pages = 100
            pages_to_process = min(len(pdf_reader.pages), max_pages)
            
            if len(pdf_reader.pages) > max_pages:
                logger.warning("[UPLOAD] PDF has %s pages, processing only first %s", len(pdf_reader.pages), max_pages)
            
            extracted_pages = []
            for i in range(pages_to_process):
                try:
                    page_text = pdf_reader.pages[i].extract_text() or ''
                    extracted_pages.append(page_text)
                except Exception as e:
                    logger.warning("[UPLOAD] Failed to extract text from page %s of %s: %s", i+1, filename, e)
                    continue
            
            extracted_text = "\n".join(extracted_pages)
            logger.info("[UPLOAD] Successfully extracted text from %s pages of PDF: %s", pages_to_process, filename)
            return extracted_text
            
    except pypdf.errors.PdfReadError as e:
        logger.error("[UPLOAD] PDF read error for %s: %s", filename, e, exc_info=True)
        raise ValueError("Invalid or corrupted PDF file")
    except Exception as e:
        logger.error("[UPLOAD] Unexpected error extracting PDF text from %s: %s", filename, e, exc_info=True)
        raise ValueError("PDF text extraction failed")

async def extract_txt_text_secure(file_path: str, filename: str) -> str:
    """Securely extract text from TXT file with encoding detection and size limits"""
    try:
        # Try multiple encodings
        encodings = ['utf-8', 'utf-16', 'latin-1', 'cp1252']
        
        for encoding in encodings:
            try:
                with open(file_path, 'r', encoding=encoding) as f:
                    # Read with size limit to prevent memory exhaustion
                    max_text_size = 10 * 1024 * 1024  # 10MB text limit
                    text = f.read(max_text_size)
                    
                    if len(text) == max_text_size:
                        logger.warning("[UPLOAD] Text file %s truncated at %s bytes", filename, max_text_size)
                    
                    logger.info("[UPLOAD] Successfully read text file %s with encoding %s", filename, encoding)
                    return text
                    
            except UnicodeDecodeError:
                continue  # Try next encoding
        
        # If all encodings failed
        raise ValueError("Unable to decode text file with any supported encoding")
        
    except Exception as e:
        logger.error("[UPLOAD] Error reading text file %s: %s", filename, e, exc_info=True)
        raise ValueError("Text file reading failed")

async def extract_csv_columns(file_path: str, filename: str) -> list:
    """Extract column names from CSV file"""
    try:
        # Read just the first row to get column names
        df = pd.read_csv(file_path, nrows=0)  # nrows=0 means read only headers
        columns = df.columns.tolist()
        
        logger.info("[UPLOAD] Extracted %s columns from CSV: %s", len(columns), filename)
        return columns
        
    except Exception as e:
        logger.warning("[UPLOAD] Failed to extract columns from CSV %s: %s", filename, e)
        return []

async def extract_excel_columns(file_path: str, filename: str) -> list:
    """Extract column names from Excel file"""
    try:
        # Read just the first row to get column names
        df = pd.read_excel(file_path, nrows=0)  # nrows=0 means read only headers
        columns = df.columns.tolist()
        
        logger.info("[UPLOAD] Extracted %s columns from Excel: %s", len(columns), filename)
        return columns
        
    except Exception as e:
        logger.warning("[UPLOAD] Failed to extract columns from Excel %s: %s", filename, e)
        return []

async def extract_excel_text_secure(file_path: str, filename: str) -> str:
    """Securely extract text from Excel files (.xlsx, .xls)"""
    if not HAS_OPENPYXL:
        logger.error("[UPLOAD] openpyxl not available for Excel text extraction: %s", filename, exc_info=True)
        raise ValueError("Excel text extraction not available - openpyxl not installed")
    
    try:
        workbook = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
        extracted_sheets = []
        
        # Limit number of sheets to prevent DoS
        max_sheets = 10
        sheets_to_process = min(len(workbook.sheetnames), max_sheets)
        
        if len(workbook.sheetnames) > max_sheets:
            logger.warning("[UPLOAD] Excel has %s sheets, processing only first %s", len(workbook.sheetnames), max_sheets)
        
        for sheet_name in workbook.sheetnames[:sheets_to_process]:
            try:
                worksheet = workbook[sheet_name]
                sheet_data = []
                
                # Limit rows to prevent memory issues
                max_rows = 1000
                row_count = 0
                
                for row in worksheet.iter_rows(values_only=True):
                    if row_count >= max_rows:
                        logger.warning("[UPLOAD] Sheet '%s' has more than %s rows, truncating", sheet_name, max_rows)
                        break
                    
                    # Convert row to string, handling None values
                    row_text = '\t'.join([str(cell) if cell is not None else '' for cell in row])
                    if row_text.strip():  # Skip empty rows
                        sheet_data.append(row_text)
                    row_count += 1
                
                if sheet_data:
                    extracted_sheets.append(f"=== Sheet: {sheet_name} ===\n" + '\n'.join(sheet_data))
                
            except Exception as e:
                logger.warning("[UPLOAD] Failed to extract data from sheet '%s' in %s: %s", sheet_name, filename, e)
                continue
        
        workbook.close()
        extracted_text = "\n\n".join(extracted_sheets)
        logger.info("[UPLOAD] Successfully extracted text from %s sheets of Excel: %s", sheets_to_process, filename)
        return extracted_text
        
    except Exception as e:
        logger.error("[UPLOAD] Error extracting Excel text from %s: %s", filename, e, exc_info=True)
        raise ValueError("Excel text extraction failed")

async def extract_docx_text_secure(file_path: str, filename: str) -> str:
    """Securely extract text from Word documents (.docx)"""
    if not HAS_PYTHON_DOCX:
        logger.error("[UPLOAD] python-docx not available for DOCX text extraction: %s", filename, exc_info=True)
        raise ValueError("DOCX text extraction not available - python-docx not installed")
    
    try:
        doc = Document(file_path)
        extracted_paragraphs = []
        
        # Limit number of paragraphs to prevent DoS
        max_paragraphs = 2000
        paragraphs_processed = 0
        
        for paragraph in doc.paragraphs:
            if paragraphs_processed >= max_paragraphs:
                logger.warning("[UPLOAD] DOCX has more than %s paragraphs, truncating", max_paragraphs)
                break
            
            text = paragraph.text.strip()
            if text:  # Skip empty paragraphs
                extracted_paragraphs.append(text)
            paragraphs_processed += 1
        
        # Also extract text from tables
        for table in doc.tables:
            table_data = []
            for row in table.rows:
                row_data = []
                for cell in row.cells:
                    cell_text = cell.text.strip()
                    if cell_text:
                        row_data.append(cell_text)
                if row_data:
                    table_data.append(' | '.join(row_data))
            
            if table_data:
                extracted_paragraphs.append("\n=== Table ===\n" + '\n'.join(table_data))
        
        extracted_text = '\n\n'.join(extracted_paragraphs)
        logger.info("[UPLOAD] Successfully extracted text from DOCX: %s", filename)
        return extracted_text
        
    except Exception as e:
        logger.error("[UPLOAD] Error extracting DOCX text from %s: %s", filename, e, exc_info=True)
        raise ValueError("DOCX text extraction failed")

async def extract_pptx_text_secure(file_path: str, filename: str) -> str:
    """Securely extract text from PowerPoint presentations (.pptx)"""
    if not HAS_PYTHON_PPTX:
        logger.error("[UPLOAD] python-pptx not available for PPTX text extraction: %s", filename, exc_info=True)
        raise ValueError("PPTX text extraction not available - python-pptx not installed")
    
    try:
        prs = Presentation(file_path)
        extracted_slides = []
        
        # Limit number of slides to prevent DoS
        max_slides = 100
        slides_to_process = min(len(prs.slides), max_slides)
        
        if len(prs.slides) > max_slides:
            logger.warning("[UPLOAD] PPTX has %s slides, processing only first %s", len(prs.slides), max_slides)
        
        for i, slide in enumerate(prs.slides[:slides_to_process]):
            try:
                slide_text = []
                
                # Extract text from all shapes in the slide
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())
                    
                    # Extract text from tables in shapes
                    if hasattr(shape, "table"):
                        table_data = []
                        for row in shape.table.rows:
                            row_data = []
                            for cell in row.cells:
                                cell_text = cell.text.strip()
                                if cell_text:
                                    row_data.append(cell_text)
                            if row_data:
                                table_data.append(' | '.join(row_data))
                        
                        if table_data:
                            slide_text.append("Table:\n" + '\n'.join(table_data))
                
                if slide_text:
                    extracted_slides.append(f"=== Slide {i+1} ===\n" + '\n\n'.join(slide_text))
                
            except Exception as e:
                logger.warning("[UPLOAD] Failed to extract text from slide %s in %s: %s", i+1, filename, e)
                continue
        
        extracted_text = "\n\n".join(extracted_slides)
        logger.info("[UPLOAD] Successfully extracted text from %s slides of PPTX: %s", slides_to_process, filename)
        return extracted_text
        
    except Exception as e:
        logger.error("[UPLOAD] Error extracting PPTX text from %s: %s", filename, e, exc_info=True)
        raise ValueError("PPTX text extraction failed")

async def extract_rtf_text_secure(file_path: str, filename: str) -> str:
    """Securely extract text from RTF files"""
    try:
        # RTF files can be read as text files, but we need to handle RTF formatting
        # For basic RTF parsing, we'll read as text and try to extract readable content
        
        # First, try to detect encoding
        encoding = 'utf-8'
        if HAS_CHARDET:
            with open(file_path, 'rb') as f:
                raw_data = f.read(10000)  # Read first 10KB for detection
                detected = chardet.detect(raw_data)
                if detected['encoding'] and detected['confidence'] > 0.7:
                    encoding = detected['encoding']
                    logger.info("[UPLOAD] Detected RTF encoding: %s (confidence: %.2f)", encoding, detected['confidence'])
        
        with open(file_path, 'r', encoding=encoding, errors='replace') as f:
            content = f.read(1024 * 1024)  # Limit to 1MB
        
        # Basic RTF content extraction (remove RTF control codes)
        # Remove RTF header and control groups
        content = re.sub(r'\\[a-z]+\d*[\s]?', ' ', content)  # Remove RTF control words
        content = re.sub(r'[{}]', '', content)  # Remove braces
        content = re.sub(r'\\[^a-z]', '', content)  # Remove other backslash sequences
        
        # Clean up whitespace
        content = re.sub(r'\s+', ' ', content)
        content = content.strip()
        
        # Extract meaningful text (filter out remaining RTF artifacts)
        lines = []
        for line in content.split('\n'):
            line = line.strip()
            # Skip lines that are mostly RTF artifacts
            if line and len(line) > 3 and not line.startswith('\\'):
                lines.append(line)
        
        extracted_text = '\n'.join(lines)
        logger.info("[UPLOAD] Successfully extracted text from RTF: %s", filename)
        return extracted_text
        
    except Exception as e:
        logger.error("[UPLOAD] Error extracting RTF text from %s: %s", filename, e, exc_info=True)
        raise ValueError("RTF text extraction failed")

@router.post("/raw-text")
async def upload_raw_text(raw_text: RawTextInput):
    """
    Accept raw text input directly without file upload.
    Useful for quick analysis of text data, SQL queries, or code snippets.
    """
    try:
        logger.info("[RAW_TEXT] Received raw text input: %s", raw_text.title)
        
        # Validate text content
        if not raw_text.text or not raw_text.text.strip():
            return {"error": "Text content cannot be empty"}
        
        # Limit text size to prevent memory issues
        max_text_size = 10 * 1024 * 1024  # 10MB
        if len(raw_text.text) > max_text_size:
            return {"error": f"Text size exceeds maximum limit of {max_text_size // (1024*1024)}MB"}
        
        # Sanitize text content
        sanitized_text = sanitize_extracted_text(raw_text.text)
        
        # Generate a filename based on title
        # Create safe filename from title
        safe_title = re.sub(r'[^\w\-_\. ]', '', raw_text.title)
        safe_title = re.sub(r'\s+', '_', safe_title).strip('_')
        if not safe_title:
            safe_title = "raw_text_input"
        

        timestamp = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
        filename = f"{safe_title}_{timestamp}.txt"
        
        # Save to uploads directory
        file_path = secure_file_path(filename)
        
        # Create the content with metadata
        content_with_metadata = f"""# {raw_text.title}
# Created: {datetime.datetime.now().isoformat()}
# Description: {raw_text.description}

{sanitized_text}"""
        
        with open(file_path, 'w', encoding='utf-8') as f:
            f.write(content_with_metadata)
        
        # Set secure permissions
        os.chmod(file_path, stat.S_IRUSR | stat.S_IWUSR)
        
        logger.info("[RAW_TEXT] Successfully saved raw text to: %s", filename)
        
        # Index the raw text into ChromaDB for RAG
        try:
            from backend.core.chromadb_client import ChromaDBClient, chunk_text
            chroma_client = ChromaDBClient()
            
            # Chunk the sanitized text
            chunks = chunk_text(sanitized_text, chunk_size=1000, overlap=200)
            logger.info("[RAW_TEXT] Chunking text into %s chunks for ChromaDB", len(chunks))
            
            # Index each chunk
            for i, chunk in enumerate(chunks):
                doc_id = f"{filename}_chunk_{i}"
                metadata = {
                    "filename": filename,
                    "chunk_index": i,
                    "total_chunks": len(chunks),
                    "source_type": "raw_text",
                    "title": raw_text.title
                }
                chroma_client.add_document(doc_id, chunk, metadata=metadata)
            
            logger.info("[RAW_TEXT] Indexed %s chunks into ChromaDB for: %s", len(chunks), filename)
        except Exception as e:
            logger.warning("[RAW_TEXT] ChromaDB indexing failed (non-critical): %s", e)
            # Don't fail the upload if ChromaDB indexing fails
        
        return {
            "filename": filename,
            "message": "Raw text uploaded successfully",
            "text_size": len(raw_text.text),
            "title": raw_text.title,
            "description": raw_text.description,
            "type": "raw_text"
        }
        
    except Exception as e:
        logger.error("[RAW_TEXT] Error processing raw text input: %s", e, exc_info=True)
        return {"error": f"Raw text processing failed: {str(e)}"}

@router.get("/preview-file/{filename}")
async def preview_file(filename: str):
    """
    Generate a preview of an uploaded file.
    Returns file content in a structured format for frontend display.
    """
    try:
        logger.info("[PREVIEW] Generating preview for file: %s", filename)
        
        # Get secure file path
        file_path = secure_file_path(filename)
        
        if not os.path.exists(file_path):
            return {"error": f"File '{filename}' not found"}
        
        # Get file metadata
        file_stats = os.stat(file_path)
        file_size = file_stats.st_size
        
        # Determine file type
        file_ext = os.path.splitext(filename)[1].lower()
        
        result = {
            "filename": filename,
            "metadata": {
                "size": file_size,
                "file_extension": file_ext
            }
        }
        
        # Generate preview based on file type
        if file_ext == '.csv':
            result = preview_csv_file(file_path, result)
        elif file_ext in ['.xlsx', '.xls']:
            result = preview_excel_file(file_path, result)
        elif file_ext == '.json':
            result = preview_json_file(file_path, result)
        elif file_ext == '.pdf':
            result = preview_pdf_file(file_path, result)
        elif file_ext in ['.txt', '.rtf', '.docx', '.pptx']:
            result = await preview_text_file(file_path, result, file_ext)
        else:
            result["error"] = f"Preview not supported for file type: {file_ext}"
        
        logger.info("[PREVIEW] Successfully generated preview for: %s", filename)
        return result
    
    except Exception as e:
        logger.error("[PREVIEW] Error generating preview for %s: %s", filename, e, exc_info=True)
        return {"error": f"Preview generation failed: {str(e)}"}

def preview_csv_file(file_path: str, result: dict) -> dict:
    """Generate preview for CSV files"""
    try:
        # Read CSV with encoding detection
        encoding = 'utf-8'
        if HAS_CHARDET:
            with open(file_path, 'rb') as f:
                raw_data = f.read(10000)
                detected = chardet.detect(raw_data)
                if detected['encoding'] and detected['confidence'] > 0.7:
                    encoding = detected['encoding']
        
        # Read CSV file
        df = pd.read_csv(file_path, encoding=encoding, nrows=100)  # Limit to first 100 rows
        df = df.where(pd.notnull(df), None)  # Replace NaN with None for JSON compliance
        
        result["data"] = df.to_dict('records')
        result["columns"] = list(df.columns)
        result["metadata"]["rows"] = len(df)
        result["metadata"]["columns"] = len(df.columns)
        result["metadata"]["encoding"] = encoding
        
        return result
    
    except Exception as e:
        result["error"] = f"CSV preview failed: {str(e)}"
        return result

def preview_excel_file(file_path: str, result: dict) -> dict:
    """Generate preview for Excel files"""
    try:
        if not HAS_OPENPYXL:
            result["error"] = "Excel preview requires openpyxl package"
            return result
        
        # Read all sheets
        excel_file = pd.ExcelFile(file_path)
        sheets_data = {}
        
        for sheet_name in excel_file.sheet_names:
            df = pd.read_excel(file_path, sheet_name=sheet_name, nrows=100)
            df = df.where(pd.notnull(df), None)  # Replace NaN with None for JSON compliance
            sheets_data[sheet_name] = {
                "data": df.to_dict('records'),
                "columns": list(df.columns),
                "rows": len(df),
                "columns_count": len(df.columns)
            }
        
        result["sheets"] = sheets_data
        result["metadata"]["sheet_count"] = len(excel_file.sheet_names)
        result["metadata"]["sheet_names"] = excel_file.sheet_names
        
        return result
    
    except Exception as e:
        result["error"] = f"Excel preview failed: {str(e)}"
        return result

def preview_json_file(file_path: str, result: dict) -> dict:
    """Generate preview for JSON files"""
    try:
        # Read JSON content
        with open(file_path, 'r', encoding='utf-8') as f:
            content = f.read(10000)  # Limit to first 10KB for preview
        
        result["content"] = content
        result["metadata"]["encoding"] = "utf-8"
        
        # Try to parse and validate JSON
        import json
        try:
            json_data = json.loads(content)
            if isinstance(json_data, list):
                result["metadata"]["json_type"] = "array"
                result["metadata"]["items"] = len(json_data)
            elif isinstance(json_data, dict):
                result["metadata"]["json_type"] = "object"
                result["metadata"]["keys"] = list(json_data.keys())[:10]  # First 10 keys
            else:
                result["metadata"]["json_type"] = "primitive"
        except json.JSONDecodeError:
            result["metadata"]["json_valid"] = False
        
        return result
    
    except Exception as e:
        result["error"] = f"JSON preview failed: {str(e)}"
        return result

def preview_pdf_file(file_path: str, result: dict) -> dict:
    """Generate preview for PDF files"""
    try:
        if not PyPDF2:
            result["error"] = "PDF preview requires PyPDF2 package"
            return result
        
        # Extract first page text
        with open(file_path, 'rb') as f:
            pdf_reader = PyPDF2.PdfReader(f)
            
            result["metadata"]["pages"] = len(pdf_reader.pages)
            
            # Extract text from first few pages (max 3)
            text_content = ""
            max_pages = min(3, len(pdf_reader.pages))
            
            for page_num in range(max_pages):
                page = pdf_reader.pages[page_num]
                text_content += page.extract_text() + "\n\n"
            
            # Limit content length
            if len(text_content) > 5000:
                text_content = text_content[:5000] + "\n... (truncated)"
            
            result["content"] = text_content
        
        return result
    
    except Exception as e:
        result["error"] = f"PDF preview failed: {str(e)}"
        return result

async def preview_text_file(file_path: str, result: dict, file_ext: str) -> dict:
    """Generate preview for text-based files"""
    try:
        encoding = 'utf-8'
        
        # Detect encoding for better text reading
        if HAS_CHARDET:
            with open(file_path, 'rb') as f:
                raw_data = f.read(10000)
                detected = chardet.detect(raw_data)
                if detected['encoding'] and detected['confidence'] > 0.7:
                    encoding = detected['encoding']
        
        # Extract content based on file type
        if file_ext == '.docx':
            content = await extract_docx_text_secure(file_path, os.path.basename(file_path))
        elif file_ext == '.pptx':
            content = await extract_pptx_text_secure(file_path, os.path.basename(file_path))
        elif file_ext == '.rtf':
            content = await extract_rtf_text_secure(file_path, os.path.basename(file_path))
        else:  # .txt and other text files
            with open(file_path, 'r', encoding=encoding, errors='replace') as f:
                content = f.read(10000)  # First 10KB
        
        # Limit content length
        if len(content) > 10000:
            content = content[:10000] + "\n... (truncated)"
        
        result["content"] = content
        result["metadata"]["encoding"] = encoding
        result["metadata"]["content_length"] = len(content)
        
        return result
    
    except Exception as e:
        result["error"] = f"Text preview failed: {str(e)}"
        return result

@router.get("/download-file/{filename}")
async def download_file(filename: str):
    """
    Download an uploaded file.
    """
    try:
        from fastapi.responses import FileResponse
        
        logger.info("[DOWNLOAD] Downloading file: %s", filename)
        
        # Get secure file path
        file_path = secure_file_path(filename)
        
        if not os.path.exists(file_path):
            return {"error": f"File '{filename}' not found"}
        
        # Return file as download
        return FileResponse(
            path=file_path,
            filename=filename,
            media_type='application/octet-stream'
        )
    
    except Exception as e:
        logger.error("[DOWNLOAD] Error downloading file %s: %s", filename, e, exc_info=True)
        return {"error": f"Download failed: {str(e)}"}
