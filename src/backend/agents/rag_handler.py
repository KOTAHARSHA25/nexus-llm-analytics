"""
RAG Handler Module
==================
Handles Retrieval-Augmented Generation (RAG) for unstructured document analysis.
Extracted from crew_manager.py for better maintainability.
"""

import logging
import os
from typing import Dict, Any, Optional, List

from backend.agents.model_initializer import get_model_initializer
from backend.core.utils import friendly_error


# Singleton instance
_rag_handler: Optional['RAGHandler'] = None


class RAGHandler:
    """
    Handles RAG-based analysis for unstructured documents (PDF, TXT, DOCX, etc.).
    Uses ChromaDB for vector storage and retrieval.
    """
    
    def __init__(self):
        self._initializer = get_model_initializer()
        self._document_store = {}
        
        logging.info("📚 RAGHandler created (lazy loading enabled)")
    
    def analyze_unstructured(
        self,
        file_path: str,
        query: str,
        enable_cot: bool = False
    ) -> Dict[str, Any]:
        """
        Analyze unstructured document using RAG.
        
        Args:
            file_path: Path to the document
            query: User's analysis query
            enable_cot: Whether to enable Chain-of-Thought reasoning
            
        Returns:
            Analysis results dictionary
        """
        try:
            if not file_path or not os.path.exists(file_path):
                return friendly_error(
                    f"Document not found: {file_path}",
                    "Please upload the document first"
                )
            
            # Extract text from document
            document_text = self._extract_text(file_path)
            
            if not document_text:
                return friendly_error(
                    "Could not extract text from document",
                    "Ensure the document is a valid PDF, TXT, or DOCX file"
                )
            
            # Generate response using LLM
            response = self._generate_rag_response(query, document_text, enable_cot)
            
            return {
                "success": True,
                "result": response,
                "filename": os.path.basename(file_path),
                "query": query,
                "type": "rag_analysis",
                "document_length": len(document_text)
            }
            
        except Exception as e:
            logging.error(f"RAG analysis failed: {e}")
            return friendly_error(f"RAG analysis failed: {str(e)}", "Check document format")
    
    def _extract_text(self, file_path: str) -> str:
        """Extract text content from various document formats."""
        file_ext = os.path.splitext(file_path)[1].lower()
        
        try:
            if file_ext == '.txt':
                with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                    return f.read()
            
            elif file_ext == '.pdf':
                return self._extract_pdf_text(file_path)
            
            elif file_ext in ['.docx', '.doc']:
                return self._extract_docx_text(file_path)
            
            elif file_ext == '.pptx':
                return self._extract_pptx_text(file_path)
            
            elif file_ext == '.rtf':
                return self._extract_rtf_text(file_path)
            
            else:
                # Try reading as plain text
                with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                    return f.read()
                    
        except Exception as e:
            logging.error(f"Text extraction failed for {file_path}: {e}")
            return ""
    
    def _extract_pdf_text(self, file_path: str) -> str:
        """Extract text from PDF files."""
        try:
            # Try pdfplumber first (better table handling)
            try:
                import pdfplumber
                text_parts = []
                with pdfplumber.open(file_path) as pdf:
                    for page in pdf.pages:
                        text = page.extract_text()
                        if text:
                            text_parts.append(text)
                return "\n\n".join(text_parts)
            except ImportError:
                pass
            
            # Fall back to PyPDF2
            try:
                from PyPDF2 import PdfReader
                reader = PdfReader(file_path)
                text_parts = []
                for page in reader.pages:
                    text = page.extract_text()
                    if text:
                        text_parts.append(text)
                return "\n\n".join(text_parts)
            except ImportError:
                pass
            
            logging.warning("No PDF library available (install pdfplumber or PyPDF2)")
            return ""
            
        except Exception as e:
            logging.error(f"PDF extraction failed: {e}")
            return ""
    
    def _extract_docx_text(self, file_path: str) -> str:
        """Extract text from DOCX files."""
        try:
            from docx import Document
            doc = Document(file_path)
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
            return "\n\n".join(paragraphs)
        except ImportError:
            logging.warning("python-docx not installed for DOCX support")
            return ""
        except Exception as e:
            logging.error(f"DOCX extraction failed: {e}")
            return ""
    
    def _extract_pptx_text(self, file_path: str) -> str:
        """Extract text from PPTX files."""
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            text_parts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text_parts.append(shape.text)
            return "\n\n".join(text_parts)
        except ImportError:
            logging.warning("python-pptx not installed for PPTX support")
            return ""
        except Exception as e:
            logging.error(f"PPTX extraction failed: {e}")
            return ""
    
    def _extract_rtf_text(self, file_path: str) -> str:
        """Extract text from RTF files."""
        try:
            from striprtf.striprtf import rtf_to_text
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                rtf_content = f.read()
            return rtf_to_text(rtf_content)
        except ImportError:
            logging.warning("striprtf not installed for RTF support")
            # Fall back to basic reading
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                return f.read()
        except Exception as e:
            logging.error(f"RTF extraction failed: {e}")
            return ""
    
    def _generate_rag_response(
        self,
        query: str,
        document_text: str,
        enable_cot: bool = False
    ) -> str:
        """Generate response using RAG approach."""
        
        # Truncate document if too long
        max_context = 4000  # tokens roughly
        if len(document_text) > max_context * 4:  # ~4 chars per token
            document_text = self._truncate_text(document_text, max_context * 4)
        
        if enable_cot:
            prompt = f"""You are analyzing a document. Use Chain-of-Thought reasoning.

DOCUMENT CONTENT:
{document_text}

QUESTION: {query}

Think step by step:
1. First, identify the relevant sections of the document
2. Extract key information related to the question
3. Analyze and synthesize the information
4. Provide a clear, comprehensive answer

ANALYSIS:"""
        else:
            prompt = f"""You are a helpful assistant analyzing a document.

DOCUMENT CONTENT:
{document_text}

QUESTION: {query}

Based on the document content above, provide a clear and comprehensive answer to the question.

ANSWER:"""
        
        try:
            response = self._initializer.llm_client.generate(
                prompt=prompt,
                adaptive_timeout=True
            )
            
            if isinstance(response, dict) and 'response' in response:
                return response['response']
            return str(response)
            
        except Exception as e:
            logging.error(f"RAG response generation failed: {e}")
            return f"Unable to generate response: {str(e)}"
    
    def _truncate_text(self, text: str, max_chars: int) -> str:
        """Truncate text to max characters, preserving sentence boundaries."""
        if len(text) <= max_chars:
            return text
        
        # Find a sentence boundary near the limit
        truncated = text[:max_chars]
        
        # Try to end at a sentence
        last_period = truncated.rfind('.')
        last_newline = truncated.rfind('\n')
        
        boundary = max(last_period, last_newline)
        if boundary > max_chars * 0.8:  # Only use if reasonably far in
            truncated = text[:boundary + 1]
        
        return truncated + "\n\n[Document truncated for processing...]"
    
    def create_embeddings(self, document_text: str, document_id: str) -> bool:
        """
        Create embeddings for document and store in ChromaDB.
        
        Args:
            document_text: Full text of the document
            document_id: Unique identifier for the document
            
        Returns:
            True if successful, False otherwise
        """
        try:
            chroma = self._initializer.chroma_client
            if not chroma:
                logging.warning("ChromaDB not available")
                return False
            
            # Split text into chunks
            chunks = self._chunk_text(document_text)
            
            # Get or create collection
            collection = chroma.get_or_create_collection(
                name="documents",
                metadata={"description": "Document embeddings for RAG"}
            )
            
            # Add chunks to collection
            collection.add(
                documents=chunks,
                ids=[f"{document_id}_{i}" for i in range(len(chunks))],
                metadatas=[{"document_id": document_id, "chunk_index": i} for i in range(len(chunks))]
            )
            
            logging.info(f"✅ Created embeddings for document {document_id}: {len(chunks)} chunks")
            return True
            
        except Exception as e:
            logging.error(f"Embedding creation failed: {e}")
            return False
    
    def _chunk_text(self, text: str, chunk_size: int = 500, overlap: int = 50) -> List[str]:
        """Split text into overlapping chunks."""
        chunks = []
        start = 0
        
        while start < len(text):
            end = start + chunk_size
            
            # Find a good breaking point
            if end < len(text):
                # Look for paragraph or sentence break
                for sep in ['\n\n', '\n', '. ', ' ']:
                    break_point = text.rfind(sep, start, end)
                    if break_point > start + chunk_size * 0.5:
                        end = break_point + len(sep)
                        break
            
            chunks.append(text[start:end])
            start = end - overlap
        
        return chunks


def get_rag_handler() -> RAGHandler:
    """Get the singleton RAGHandler instance."""
    global _rag_handler
    if _rag_handler is None:
        _rag_handler = RAGHandler()
    return _rag_handler
